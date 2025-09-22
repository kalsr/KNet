# python
import streamlit as st
import pandas as pd
import numpy as np
import yfinance as yf
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
import io, datetime, base64

# Try to import python-pptx safely
try:
    from pptx import Presentation
    from pptx.util import Inches
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# ---------------------------
# Page config
# ---------------------------
st.set_page_config(page_title="Stock Comparison Tool - Designed By Randy Singh", layout="wide")

# ---------------------------
# Helper functions
# ---------------------------
def fetch_company_data(ticker):
    tk = yf.Ticker(ticker)
    info = tk.info
    hist = tk.history(period='10y', interval='1d', auto_adjust=False)
    return {'info': info, 'history': hist}

def compute_simple_metrics(data):
    info = data['info']
    price = info.get('regularMarketPrice') or info.get('previousClose')
    div_yield = info.get('dividendYield', None)
    if div_yield is not None:
        div_yield *= 100
    pe = info.get('trailingPE', None)
    de_ratio = info.get('debtToEquity', None)
    revenue = info.get('totalRevenue', None)
    net_income = info.get('netIncomeToCommon', None) or info.get('netIncome', None)
    rev_bil = revenue / 1e9 if revenue else None
    ni_bil = net_income / 1e9 if net_income else None
    return {'price': price, 'dividend_yield_pct': div_yield, 'pe': pe,
            'de_ratio': de_ratio, 'revenue_bil': rev_bil, 'net_income_bil': ni_bil}

def minmax_norm(series, invert=False):
    s = series.dropna()
    if s.empty:
        return pd.Series([0.5]*len(series), index=series.index)
    mn, mx = s.min(), s.max()
    if mx == mn:
        return pd.Series([0.5]*len(series), index=series.index)
    norm = (series - mn) / (mx - mn)
    return 1.0 - norm if invert else norm

def one_year_return(hist):
    if hist is None or hist.empty: return None
    closes = hist['Close'].dropna()
    if len(closes) < 252: return None
    return (closes.iloc[-1] - closes.iloc[-252]) / closes.iloc[-252] * 100

# ---------------------------
# UI Header
# ---------------------------
st.markdown("<h1 style='color:#0b6fb0'> Stock Comparison Tool</h1>", unsafe_allow_html=True)
st.markdown("Compare any two companies using Yahoo Finance data. Includes ratios, price history, and report exports. **Not financial advice.**")

# ---------------------------
# Sidebar Ticker Selection
# ---------------------------
popular_tickers = [
    "AAPL","MSFT","AMZN","GOOGL","META","TSLA","NVDA","NFLX",
    "KO","PEP","JPM","BAC","V","MA","DIS","WMT","PFE","JNJ","XOM","CVX"
]

st.sidebar.header("Select Companies")
ticker1 = st.sidebar.selectbox("Company 1", popular_tickers, index=0)
ticker2 = st.sidebar.selectbox("Company 2", popular_tickers, index=1)

manual1 = st.sidebar.text_input("Or enter custom ticker for Company 1", "")
manual2 = st.sidebar.text_input("Or enter custom ticker for Company 2", "")

if manual1.strip():
    ticker1 = manual1.strip().upper()
if manual2.strip():
    ticker2 = manual2.strip().upper()

# ---------------------------
# Fetch Data
# ---------------------------
tickers = {ticker1: ticker1, ticker2: ticker2}
data, metrics = {}, []
for t in tickers:
    try:
        data[t] = fetch_company_data(t)
        m = compute_simple_metrics(data[t])
        metrics.append({
            'company': tickers[t],
            'ticker': t,
            'price': m['price'],
            'revenue_bil': m['revenue_bil'],
            'net_income_bil': m['net_income_bil'],
            'dividend_yield_pct': m['dividend_yield_pct'],
            'pe': m['pe'],
            'de_ratio': m['de_ratio'],
            '1y_return_pct': one_year_return(data[t]['history'])
        })
    except Exception as e:
        st.error(f"Failed to fetch {t}: {e}")

df = pd.DataFrame(metrics)
if not df.empty:
    df['net_margin'] = df['net_income_bil'] / df['revenue_bil']
    df['earnings_yield'] = df['pe'].apply(lambda x: 1/x if pd.notnull(x) and x!=0 else np.nan)

st.subheader(" Key Metrics")
if not df.empty:
    st.dataframe(df.set_index('company'))
else:
    st.warning("No data to show.")

# ---------------------------
# Charts
# ---------------------------
if not df.empty:
    st.subheader(" 10-Year Price History")
    fig, ax = plt.subplots(figsize=(10,4))
    for t in tickers:
        ax.plot(data[t]['history'].index, data[t]['history']['Close'], label=t)
    ax.legend(); ax.set_title("Price Comparison"); ax.set_ylabel("USD")
    st.pyplot(fig)

# ---------------------------
# PDF Export
# ---------------------------
if not df.empty and st.button(" Download PDF Summary"):
    out_buf = io.BytesIO()
    with PdfPages(out_buf) as pdf:
        fig, ax = plt.subplots(figsize=(8.5,11))
        ax.axis('off')
        txt = "Stock Comparison Report\n\n"
        for _, row in df.iterrows():
            txt += f"{row['company']} ({row['ticker']}): Price ${row['price']}, Revenue {row['revenue_bil']}B, Net Margin {row['net_margin']:.2f}, DivYield {row['dividend_yield_pct']}%, P/E {row['pe']}\n"
        ax.text(0.01, 0.99, txt, va='top', fontsize=10, family='monospace')
        pdf.savefig(fig); plt.close(fig)
    out_buf.seek(0)
    b64 = base64.b64encode(out_buf.read()).decode()
    href = f'<a href="data:application/pdf;base64,{b64}" download="Stock_Comparison_Report.pdf"><button>Download PDF</button></a>'
    st.markdown(href, unsafe_allow_html=True)

# ---------------------------
# PPTX Export
# ---------------------------
if not df.empty:
    if PPTX_AVAILABLE:
        if st.button(" Download PPTX Summary"):
            prs = Presentation()
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = "Stock Comparison Report"
            slide.placeholders[1].text = f"Generated {datetime.datetime.utcnow().strftime('%Y-%m-%d %H:%M:%S UTC')}"

            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)
            shapes = slide.shapes
            rows, cols = df.shape[0]+1, df.shape[1]
            left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(3)
            table = shapes.add_table(rows, cols, left, top, width, height).table

            for i, col in enumerate(df.columns):
                table.cell(0,i).text = col
            for r in range(df.shape[0]):
                for c in range(df.shape[1]):
                    val = df.iloc[r,c]
                    table.cell(r+1,c).text = str(round(val,2) if isinstance(val,(int,float)) else val)

            out_buf = io.BytesIO()
            prs.save(out_buf)
            out_buf.seek(0)
            b64 = base64.b64encode(out_buf.read()).decode()
            href = f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64}" download="Stock_Comparison_Report.pptx"><button>Download PPTX</button></a>'
            st.markdown(href, unsafe_allow_html=True)
    else:
        st.info(" PPTX export requires `python-pptx`. Add it to requirements.txt: `python-pptx`")

st.markdown("---")
st.caption("This tool is designed by Randy Singh from KNet Consulting & uses Yahoo Finance data. Educational use only. Not financial advice.")
