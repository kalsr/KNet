import streamlit as st
import pandas as pd
import networkx as nx
from pyvis.network import Network
import plotly.express as px
import json
import base64
from io import BytesIO

from reportlab.platypus import SimpleDocTemplate, Paragraph
from reportlab.lib.styles import getSampleStyleSheet
from docx import Document
from pptx import Presentation

# -----------------------------
# PAGE CONFIG
# -----------------------------
st.set_page_config(page_title="KNet Cloud Computing Framework", layout="wide")

# -----------------------------
# TOP BLUE HEADER BAR (FIXED)
# -----------------------------
st.markdown("""
<div style="
background-color:#0B3D91;
padding:18px;
border-radius:10px;
text-align:center;
color:white;
font-size:26px;
font-weight:bold;">
KNet Cloud Computing Framework
</div>
<div style="
background-color:#0B3D91;
padding:10px;
border-radius:10px;
text-align:center;
color:white;
font-size:14px;
font-weight:bold;
margin-top:5px;">
Enterprise Cloud Taxonomy | MITRE | Zero Trust | AI Engine
</div>

<div style="
text-align:center;
color:#0B3D91;
font-size:14px;
font-weight:bold;
margin-top:5px;">
Developed by Randy Singh from Kalsnet (KNet) Consulting Group
</div>
""", unsafe_allow_html=True)

# -----------------------------
# SESSION STATE
# -----------------------------
if "nodes" not in st.session_state:
    st.session_state.nodes = []

# -----------------------------
# SIDEBAR - CLOUD BUILDER
# -----------------------------
st.sidebar.header(" Cloud Builder Node Input")

st.sidebar.markdown("""
**Node Name:** Logical identifier for cloud component  
Example: `Auth-Service`, `Storage-Layer`, `Web-App`
""")

node_name = st.sidebar.text_input("Node Name")

st.sidebar.markdown("""
**Cloud Type Explanation:**
- IaaS → Infrastructure (VMs, storage, networks)
- PaaS → Application runtime (APIs, Kubernetes, middleware)
- SaaS → End-user applications (email, CRM, Office tools)
""")

cloud_type = st.sidebar.selectbox("Cloud Type", ["IaaS", "PaaS", "SaaS"])

st.sidebar.markdown("""
**MITRE Techniques:**
- T1078 → Valid accounts (stolen credentials)
- T1190 → Exploit public-facing application
- T1059 → Command execution via scripts
""")

mitre = st.sidebar.selectbox(
    "MITRE Technique",
    ["T1078 - Valid Accounts", "T1190 - Exploit Public App", "T1059 - Command Execution"]
)

risk = st.sidebar.slider("Risk Score (0-100)", 0, 100, 50)

st.sidebar.markdown("""
**Zero Trust Layers:**
Identity → Device → Network → Application → Data
""")

zero_trust = st.sidebar.selectbox(
    "Zero Trust Layer",
    ["Identity", "Device", "Network", "Application", "Data"]
)

# -----------------------------
# ADD NODE (MULTI NODE SUPPORT FIX)
# -----------------------------
if st.sidebar.button("➕ Add Cloud Node"):
    if node_name:
        st.session_state.nodes.append({
            "name": node_name,
            "type": cloud_type,
            "mitre": mitre,
            "risk": risk,
            "zt": zero_trust
        })
        st.success(f"Node '{node_name}' added successfully!")

if st.sidebar.button(" Reset All Nodes"):
    st.session_state.nodes = []

# -----------------------------
# DATAFRAME
# -----------------------------
df = pd.DataFrame(st.session_state.nodes)

st.subheader(" Cloud Nodes Overview")
if not df.empty:
    st.dataframe(df, use_container_width=True)

# -----------------------------
# AI RECOMMENDATION ENGINE (EXPLAINED)
# -----------------------------
def ai_recommend(nodes):
    """
    AI engine logic:
    1. Aggregates risk scores
    2. Detects system-wide exposure
    3. Suggests architecture changes
    """

    if not nodes:
        return "No cloud nodes detected."

    avg_risk = sum(n["risk"] for n in nodes) / len(nodes)

    if avg_risk > 70:
        return """
 HIGH RISK DETECTED  
AI Recommendation:
- Enforce Zero Trust everywhere
- Migrate critical workloads to PaaS
- Enable IAM + MFA + micro-segmentation
"""
    elif avg_risk > 40:
        return """
 MEDIUM RISK  
AI Recommendation:
- Strengthen identity layer
- Add network segmentation
- Reduce exposed SaaS services
"""
    else:
        return """
 LOW RISK  
AI Recommendation:
- Optimize SaaS usage
- Consolidate workloads
- Improve monitoring
"""

st.subheader(" AI Recommendation Engine (How it works)")
st.info(ai_recommend(st.session_state.nodes))

st.markdown("""
**How AI Engine Works:**
- Collects all node risk scores
- Calculates average exposure
- Maps MITRE + Zero Trust + Cloud type
- Generates risk-based architecture guidance
""")

# -----------------------------
# NETWORK GRAPH
# -----------------------------
st.subheader(" Cloud Architecture Graph")

def build_graph(data):
    G = nx.Graph()

    for n in data:
        G.add_node(n["name"])
        G.add_edge(n["name"], n["type"])
        G.add_edge(n["name"], n["mitre"])
        G.add_edge(n["name"], n["zt"])

    return G

if st.button("📡 Generate Network Graph"):
    if df.empty:
        st.warning("No nodes available")
    else:
        G = build_graph(st.session_state.nodes)

        net = Network(height="600px", width="100%", bgcolor="#0b1a2f", font_color="white")

        for node in G.nodes:
            net.add_node(node)

        for edge in G.edges:
            net.add_edge(edge[0], edge[1])

        net.save_graph("graph.html")

        with open("graph.html", "r", encoding="utf-8") as f:
            st.components.v1.html(f.read(), height=600)

# -----------------------------
# MITRE + ZERO TRUST VISUALIZATION
# -----------------------------
st.subheader(" MITRE + Zero Trust Distribution")

if not df.empty:
    fig = px.histogram(df, x="zt", color="type",
                       title="Zero Trust Layer Distribution Across Cloud Types")
    st.plotly_chart(fig, use_container_width=True)

# -----------------------------
# DRAG & DROP EXPLANATION (REAL CLARITY)
# -----------------------------
st.subheader(" Cloud Builder Explanation")

st.markdown("""
### What is Cloud Builder Node?

A **Node** represents a cloud component:

Examples:
- Auth-Service (Identity layer)
- Storage Cluster (IaaS)
- API Gateway (PaaS)
- Email System (SaaS)

You can add MULTIPLE nodes to simulate real architectures.

Each node connects to:
- Cloud Type (IaaS / PaaS / SaaS)
- MITRE Attack Mapping
- Zero Trust Security Layer
""")

# -----------------------------
# EXPORT HELPERS (FIXED DOWNLOAD)
# -----------------------------
def to_pdf(data):
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer)
    styles = getSampleStyleSheet()
    text = Paragraph(str(data), styles["Normal"])
    doc.build([text])
    buffer.seek(0)
    return buffer

def to_word(data):
    doc = Document()
    doc.add_heading("Cloud Computing Framework Report", 0)
    doc.add_paragraph(str(data))
    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

def to_ppt(data):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Cloud Computing Framework Report"
    slide.placeholders[1].text = str(data)

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

# -----------------------------
# EXPORT SECTION (FIXED DOWNLOAD BUTTONS)
# -----------------------------
st.subheader(" Export Reports (Fixed)")

pdf_data = to_pdf(st.session_state.nodes)
st.download_button("⬇ Download PDF", pdf_data, file_name="cloud_report.pdf")

word_data = to_word(st.session_state.nodes)
st.download_button("⬇ Download Word", word_data, file_name="cloud_report.docx")

ppt_data = to_ppt(st.session_state.nodes)
st.download_button("⬇ Download PowerPoint", ppt_data, file_name="cloud_report.pptx")

json_data = json.dumps(st.session_state.nodes, indent=2)
st.download_button("⬇ Download JSON", json_data, file_name="cloud_data.json")
