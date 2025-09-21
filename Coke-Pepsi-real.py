

# coke_pepsi_real.py
import streamlit as st
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
import io, datetime, base64, textwrap, json, sqlite3, os
from pathlib import Path

# ---------- Try import yfinance gracefully ----------
try:
    import yfinance as yf
except Exception:
    st.set_page_config(page_title="Coke vs Pepsi — Pro Analyzer", layout="wide")
    st.error(
        "Missing Python package `yfinance`. Install with `pip install yfinance` and restart the app.\n"
        "If you are on Streamlit Cloud, add `yfinance` to requirements.txt."
    )
    st.stop()

# ---------- Try import python-pptx gracefully ----------
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except Exception:
    PPTX_AVAILABLE = False
else:
    PPTX_AVAILABLE = True

# ---------- Page config and styling ----------
st.set_page_config(page_title="Coke vs Pepsi — Pro Analyzer", layout="wide")
st.markdown("""
    <style>
      .stApp { background-color: #f6f9fb; }
      h1 { color: #0b6fb0; }
      .section-card { background: #ffffff; padding: 16px; border-radius: 10px; box-shadow: 0 6px 18px rgba(13,40,68,0.08); }
      .small-muted { color: #6c757d; font-size: 0.9rem; }
      .metric-card { background: #e9f3fb; padding: 10px; border-radius: 8px; }
    </style>
""", unsafe_allow_html=True)

# ---------- SQLite snapshot cache (local file) ----------
DB_PATH = "ticker_snapshots.db"

def init_db():
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute("""
        CREATE TABLE IF NOT EXISTS snapshots (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            ticker TEXT NOT NULL,
            ts TEXT NOT NULL,
            info_json TEXT,
            history_csv BLOB
        )
    """)
    conn.commit()
    conn.close()

def save_snapshot(ticker, info_dict, history_df):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    ts = datetime.datetime.utcnow().isoformat()
    info_json = json.dumps(info_dict, default=str)
    history_csv = history_df.to_csv(index=True).encode('utf-8') if history_df is not None else None
    c.execute("INSERT INTO snapshots (ticker, ts, info_json, history_csv) VALUES (?, ?, ?, ?)",
              (ticker.upper(), ts, info_json, history_csv))
    conn.commit()
    conn.close()
    return ts

def get_snapshots_for_ticker(ticker, limit=10):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute("SELECT id, ticker, ts, info_json FROM snapshots WHERE ticker=? ORDER BY id DESC LIMIT ?", (ticker.upper(), limit))
    rows = c.fetchall()
    conn.close()
    return rows

def load_snapshot_history_csv(snapshot_id):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute("SELECT history_csv FROM snapshots WHERE id=?", (snapshot_id,))
    row = c.fetchone()
    conn.close()
    if row and row[0]:
        return pd.read_csv(io.BytesIO(row[0]), index_col=0, parse_dates=True)
    return None

init_db()

# ---------- Helper functions ----------
@st.cache_data(ttl=3600)
def fetch_ticker_full(ticker_symbol, period='5y', interval='1d'):
    tk = yf.Ticker(ticker_symbol)
    info = tk.info or {}
    try:
        hist = tk.history(period=period, interval=interval, auto_adjust=False)
    except Exception:
        hist = pd.DataFrame()
    # quarterly & annual financials
    try:
        cashflow = tk.cashflow
    except Exception:
        cashflow = None
    try:
        fin = tk.financials
    except Exception:
        fin = None
    try:
        bs = tk.balance_sheet
    except Exception:
        bs = None
    try:
        q_fin = tk.quarterly_financials
    except Exception:
        q_fin = None
    try:
        q_cash = tk.quarterly_cashflow
    except Exception:
        q_cash = None
    return {
        'info': info,
        'history': hist,
        'cashflow': cashflow,
        'financials': fin,
        'balance_sheet': bs,
        'q_financials': q_fin,
        'q_cashflow': q_cash
    }

def safe_div(a,b):
    try:
        return a/b
    except Exception:
        return np.nan

def extract_many_metrics(info):
    m = {}
    m['shortName'] = info.get('shortName')
    m['marketCap'] = info.get('marketCap')
    m['enterpriseValue'] = info.get('enterpriseValue')
    m['regularMarketPrice'] = info.get('regularMarketPrice') or info.get('previousClose')
    m['previousClose'] = info.get('previousClose')
    m['sharesOutstanding'] = info.get('sharesOutstanding')
    m['trailingPE'] = info.get('trailingPE')
    m['forwardPE'] = info.get('forwardPE')
    m['pegRatio'] = info.get('pegRatio')
    m['beta'] = info.get('beta')
    m['dividendYield'] = info.get('dividendYield')
    m['trailingAnnualDividendRate'] = info.get('trailingAnnualDividendRate')
    m['totalRevenue'] = info.get('totalRevenue')
    m['totalCash'] = info.get('totalCash')
    m['totalDebt'] = info.get('totalDebt') or info.get('debt')
    m['ebitda'] = info.get('ebitda')
    m['grossMargins'] = info.get('grossMargins')
    m['operatingMargins'] = info.get('operatingMargins')
    m['profitMargins'] = info.get('profitMargins')
    m['returnOnEquity'] = info.get('returnOnEquity')
    m['returnOnAssets'] = info.get('returnOnAssets')
    m['currentRatio'] = info.get('currentRatio')
    m['epsTrailingTwelveMonths'] = info.get('epsTrailingTwelveMonths')
    # derived readable
    m['dividendYieldPct'] = m['dividendYield'] * 100.0 if m.get('dividendYield') else None
    m['totalRevenueB'] = safe_div(m.get('totalRevenue'), 1e9) if m.get('totalRevenue') else None
    m['ebitdaB'] = safe_div(m.get('ebitda'), 1e9) if m.get('ebitda') else None
    return m

def extract_fcf_history(cashflow_df):
    if cashflow_df is None or cashflow_df.empty:
        return None
    idx = [str(i).lower() for i in cashflow_df.index]
    for candidate in ['free cash flow', 'freecashflow', 'free cashflows', 'freecashflows']:
        for j,name in enumerate(idx):
            if candidate in name:
                return cashflow_df.iloc[j].fillna(0).astype(float).values.tolist()
    ocf, capex = None, None
    for j,name in enumerate(idx):
        if 'operat' in name and 'cash' in name:
            ocf = cashflow_df.iloc[j].astype(float)
        if 'capital' in name and 'expend' in name:
            capex = cashflow_df.iloc[j].astype(float)
    if ocf is not None and capex is not None:
        return (ocf - capex).fillna(0).astype(float).values.tolist()
    return None

def dcf_5yr_simple(base_fcf, growth_rate, terminal_growth, discount_rate, shares_outstanding):
    if base_fcf is None:
        return None
    years = np.arange(1,6)
    projected = [base_fcf * ((1 + growth_rate) ** y) for y in years]
    terminal_fcf = projected[-1] * (1 + terminal_growth)
    terminal_value = terminal_fcf / (discount_rate - terminal_growth) if discount_rate > terminal_growth else np.nan
    discounted = [projected[i] / ((1 + discount_rate) ** (i+1)) for i in range(5)]
    discounted_terminal = terminal_value / ((1 + discount_rate) ** 5) if not np.isnan(terminal_value) else np.nan
    pv_total = np.nansum(discounted) + (discounted_terminal if not np.isnan(discounted_terminal) else 0.0)
    per_share = pv_total / shares_outstanding if shares_outstanding and shares_outstanding > 0 else pv_total
    return {
        'pv_total': pv_total,
        'intrinsic_per_share': per_share,
        'projected': projected,
        'terminal_value': terminal_value,
        'discounted': discounted,
        'discounted_terminal': discounted_terminal
    }

def one_year_return_pct(hist_df):
    if hist_df is None or hist_df.empty:
        return np.nan
    closes = hist_df['Close'].dropna()
    if len(closes) < 2:
        return np.nan
    return (closes.iloc[-1] - closes.iloc[0]) / closes.iloc[0] * 100.0

def df_to_csv_bytes(df):
    return df.to_csv(index=True).encode('utf-8')

def make_pdf_report(analysis_summary_text, combined_df, price_histories, sens_results=None):
    buf = io.BytesIO()
    with PdfPages(buf) as pdf:
        fig, ax = plt.subplots(figsize=(8.5,11))
        ax.axis('off')
        ax.text(0, 1, analysis_summary_text, va='top', wrap=True, fontsize=10, family='monospace')
        pdf.savefig(fig); plt.close(fig)

        fig2, ax2 = plt.subplots(figsize=(11,6))
        ax2.axis('off')
        ax2.table(cellText=combined_df.round(4).values, colLabels=combined_df.columns, rowLabels=combined_df.index, loc='center')
        pdf.savefig(fig2); plt.close(fig2)

        fig3, ax3 = plt.subplots(figsize=(11,4))
        for label, hist in price_histories.items():
            if hist is not None and not hist.empty:
                ax3.plot(hist.index, hist['Close'], label=label)
        ax3.set_title("Price History")
        ax3.legend()
        pdf.savefig(fig3); plt.close(fig3)

        if sens_results:
            fig4, ax4 = plt.subplots(figsize=(6,6))
            ax4.pie([sens_results['a_pct'], sens_results['b_pct']],
                   labels=[f"{list(price_histories.keys())[0]} {sens_results['a_pct']:.1f}%", f"{list(price_histories.keys())[1]} {sens_results['b_pct']:.1f}%"],
                   startangle=140, autopct='%1.1f%%')
            ax4.set_title("Sensitivity robustness")
            pdf.savefig(fig4); plt.close(fig4)
    buf.seek(0)
    return buf

def make_pptx_report(title_text, combined_df, dcf_results, price_histories, out_path="/tmp/coke_pepsi_summary.pptx"):
    prs = Presentation()
    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Coke vs Pepsi — Analysis Summary"
    subtitle = slide.placeholders[1]
    subtitle.text = title_text

    # Metrics slide
    slide_layout = prs.slide_layouts[5]  # blank-ish
    slide = prs.slides.add_slide(slide_layout)
    left = Inches(0.5); top = Inches(1); width = Inches(9); height = Inches(5)
    tx = slide.shapes.add_textbox(left, top, width, height).text_frame
    tx.text = "Key comparative metrics"
    tx.add_paragraph()
    # add simple metrics lines
    for idx, row in combined_df.iterrows():
        tx.add_paragraph(f"{idx}: Price={row.get('Price')}, MarketCap={row.get('MarketCap')}, P/E={row.get('P/E')}")

    # DCF slide
    slide = prs.slides.add_slide(slide_layout)
    tx2 = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(5)).text_frame
    tx2.text = "DCF results (simple template):"
    for k,v in dcf_results.items():
        if v:
            tx2.add_paragraph(f"{k}: Intrinsic/share = {v['intrinsic_per_share']:.2f}, PV total = {v['pv_total']:.2f}")
        else:
            tx2.add_paragraph(f"{k}: DCF not available")

    # Save pptx
    prs.save(out_path)
    return out_path

# ---------- UI ----------
st.title("Coke vs Pepsi — Pro Financial Comparator")
st.markdown("This app fetches real-world company data (prices, fundamentals, cashflows) and creates reports, sensitivity heatmaps, and a downloadable PPTX. It stores snapshots in a local SQLite DB for reproducibility.")

# Sidebar controls
with st.sidebar:
    st.header("Controls")
    period = st.selectbox("Historical range", options=['1y','2y','5y','10y','max'], index=2)
    interval = st.selectbox("Price interval", options=['1d','1wk','1mo'], index=0)
    discount_rate = st.slider("Discount rate (decimal)", 0.01, 0.20, 0.09, 0.005, format="%.3f")
    terminal_growth = st.slider("Terminal growth (decimal)", 0.0, 0.05, 0.02, 0.001, format="%.3f")
    default_fc_growth = st.slider("Default 5yr growth (decimal)", -0.05, 0.30, 0.05, 0.005, format="%.3f")
    fetch_button = st.button(" Fetch / Refresh (and snapshot)")
    st.markdown("---")
    st.header("Sensitivity / Sampling")
    n_samples = st.number_input("Weight samples", min_value=200, max_value=5000, value=2000, step=100)
    run_sens = st.button("Run weight-sampling sensitivity")
    st.markdown("---")
    st.header("Exports")
    pptx_button = st.button(" Create PPTX Summary")
    pdf_button = st.button(" Create PDF Report")

# Tickers input
col1, col2 = st.columns(2)
with col1:
    ticker_a = st.text_input("Ticker A (left)", value="KO")
with col2:
    ticker_b = st.text_input("Ticker B (right)", value="PEP")

# fetch data (use cached fetch_ticker_full)
with st.spinner("Fetching data... (cached)"):
    tA = fetch_ticker_full(ticker_a, period=period, interval=interval)
    tB = fetch_ticker_full(ticker_b, period=period, interval=interval)

# snapshot if requested
if fetch_button:
    try:
        save_snapshot(ticker_a, tA['info'], tA['history'])
        save_snapshot(ticker_b, tB['info'], tB['history'])
        st.success("Snapshots saved to local SQLite DB")
    except Exception as e:
        st.error(f"Snapshot save failed: {e}")

# extract metrics
mA = extract_many_metrics(tA['info'])
mB = extract_many_metrics(tB['info'])

# top metrics display
st.subheader("Top-line comparison")
cA, cB = st.columns(2)
with cA:
    st.markdown(f"### {mA.get('shortName','')} ({ticker_a.upper()})")
    st.metric("Price", f"${mA.get('regularMarketPrice', np.nan):,.2f}")
    st.write(f"Market Cap: {mA.get('marketCap'):,}" if mA.get('marketCap') else "Market Cap: N/A")
    st.write(f"P/E (trailing): {mA.get('trailingPE')}" if mA.get('trailingPE') else "P/E: N/A")
    st.write(f"Dividend yield: {mA.get('dividendYieldPct'):.2f}%" if mA.get('dividendYieldPct') else "Div Yield: N/A")
with cB:
    st.markdown(f"### {mB.get('shortName','')} ({ticker_b.upper()})")
    st.metric("Price", f"${mB.get('regularMarketPrice', np.nan):,.2f}")
    st.write(f"Market Cap: {mB.get('marketCap'):,}" if mB.get('marketCap') else "Market Cap: N/A")
    st.write(f"P/E (trailing): {mB.get('trailingPE')}" if mB.get('trailingPE') else "P/E: N/A")
    st.write(f"Dividend yield: {mB.get('dividendYieldPct'):.2f}%" if mB.get('dividendYieldPct') else "Div Yield: N/A")

# combined dataframe
combined_df = pd.DataFrame({
    ticker_a.upper(): {
        'Price': mA.get('regularMarketPrice'),
        'MarketCap': mA.get('marketCap'),
        'EV': mA.get('enterpriseValue'),
        'P/E': mA.get('trailingPE'),
        'PEG': mA.get('pegRatio'),
        'EBITDA (B)': mA.get('ebitdaB'),
        'Revenue (B)': mA.get('totalRevenueB'),
        'DivYield %': mA.get('dividendYieldPct'),
        'ROE': mA.get('returnOnEquity'),
        'SharesOutstanding': mA.get('sharesOutstanding')
    },
    ticker_b.upper(): {
        'Price': mB.get('regularMarketPrice'),
        'MarketCap': mB.get('marketCap'),
        'EV': mB.get('enterpriseValue'),
        'P/E': mB.get('trailingPE'),
        'PEG': mB.get('pegRatio'),
        'EBITDA (B)': mB.get('ebitdaB'),
        'Revenue (B)': mB.get('totalRevenueB'),
        'DivYield %': mB.get('dividendYieldPct'),
        'ROE': mB.get('returnOnEquity'),
        'SharesOutstanding': mB.get('sharesOutstanding')
    }
}).T

st.subheader("Combined quick table")
st.dataframe(combined_df)

# show quarterly & TTM data (if available)
st.subheader("Quarterly & TTM snapshot (if available)")
qcols = st.columns(2)
with qcols[0]:
    st.markdown(f"#### {ticker_a.upper()} quarterly financials (head)")
    if tA['q_financials'] is not None:
        st.dataframe(tA['q_financials'].head().fillna(''))
    else:
        st.markdown("Quarterly financials not available via yfinance for this ticker.")
with qcols[1]:
    st.markdown(f"#### {ticker_b.upper()} quarterly financials (head)")
    if tB['q_financials'] is not None:
        st.dataframe(tB['q_financials'].head().fillna(''))
    else:
        st.markdown("Quarterly financials not available via yfinance for this ticker.")

# compute FCF history and DCF
fcfA = extract_fcf_history(tA['cashflow'])
fcfB = extract_fcf_history(tB['cashflow'])
baseA = np.mean(fcfA[:3]) if fcfA else (mA.get('totalRevenueB')*1e9*0.10 if mA.get('totalRevenueB') else None)
baseB = np.mean(fcfB[:3]) if fcfB else (mB.get('totalRevenueB')*1e9*0.10 if mB.get('totalRevenueB') else None)
sharesA = mA.get('sharesOutstanding') or 1.0
sharesB = mB.get('sharesOutstanding') or 1.0

# user tuned growth inputs
st.subheader("DCF assumptions & results")
gA = st.number_input(f"{ticker_a} 5yr FCF growth (decimal)", value=float(default_fc_growth), step=0.005, format="%.3f", key='gA')
gB = st.number_input(f"{ticker_b} 5yr FCF growth (decimal)", value=float(default_fc_growth), step=0.005, format="%.3f", key='gB')
tgA = st.number_input(f"{ticker_a} terminal growth (decimal)", value=float(terminal_growth), step=0.001, format="%.3f", key='tgA')
tgB = st.number_input(f"{ticker_b} terminal growth (decimal)", value=float(terminal_growth), step=0.001, format="%.3f", key='tgB')

dcfA = dcf_5yr_simple(baseA, gA, tgA, discount_rate, sharesA)
dcfB = dcf_5yr_simple(baseB, gB, tgB, discount_rate, sharesB)

dcf_show = pd.DataFrame({
    ticker_a.upper(): {'Intrinsic/share': dcfA['intrinsic_per_share'] if dcfA else np.nan, 'PV total': dcfA['pv_total'] if dcfA else np.nan},
    ticker_b.upper(): {'Intrinsic/share': dcfB['intrinsic_per_share'] if dcfB else np.nan, 'PV total': dcfB['pv_total'] if dcfB else np.nan}
}).T
st.dataframe(dcf_show)

# Composite scoring
st.subheader("Composite scoring & weight sliders")
w_pe = st.slider("Weight P/E", 0.0, 1.0, 0.25)
w_div = st.slider("Weight Dividend Yield", 0.0, 1.0, 0.25)
w_dcf = st.slider("Weight DCF intrinsic value", 0.0, 1.0, 0.3)
w_roe = st.slider("Weight ROE", 0.0, 1.0, 0.2)
total_w = (w_pe + w_div + w_dcf + w_roe) or 1.0
weights = {'pe': w_pe/total_w, 'div': w_div/total_w, 'dcf': w_dcf/total_w, 'roe': w_roe/total_w}

score_df = pd.DataFrame(index=[ticker_a.upper(), ticker_b.upper()])
score_df['pe'] = [mA.get('trailingPE'), mB.get('trailingPE')]
score_df['pe_norm'] = minmax_norm(score_df['pe'], invert=True)
score_df['div'] = [mA.get('dividendYieldPct'), mB.get('dividendYieldPct')]
score_df['div_norm'] = minmax_norm(score_df['div'])
score_df['dcf'] = [dcfA['intrinsic_per_share'] if dcfA else np.nan, dcfB['intrinsic_per_share'] if dcfB else np.nan]
score_df['dcf_norm'] = minmax_norm(score_df['dcf'])
score_df['roe'] = [mA.get('returnOnEquity'), mB.get('returnOnEquity')]
score_df['roe_norm'] = minmax_norm(score_df['roe'])

score_df['composite'] = (score_df['pe_norm']*weights['pe'] + score_df['div_norm']*weights['div'] + score_df['dcf_norm']*weights['dcf'] + score_df['roe_norm']*weights['roe'])
st.dataframe(score_df[['pe','div','dcf','roe','composite']].round(4))

# Sensitivity sampling
sens_results = None
if run_sens:
    samples = int(n_samples)
    rng = np.random.default_rng(123456)
    wins = {ticker_a.upper():0, ticker_b.upper():0}
    for i in range(samples):
        w = rng.dirichlet(alpha=np.ones(4))
        scores = score_df['pe_norm']*w[0] + score_df['div_norm']*w[1] + score_df['dcf_norm']*w[2] + score_df['roe_norm']*w[3]
        winner = scores.idxmax()
        wins[winner] += 1
    a_pct = wins[ticker_a.upper()]/samples*100.0
    b_pct = wins[ticker_b.upper()]/samples*100.0
    sens_results = {'wins':wins, 'a_pct':a_pct, 'b_pct':b_pct, 'samples':samples}
    st.success("Sensitivity sampling complete")
    figp, axp = plt.subplots(figsize=(4,3))
    axp.pie([a_pct,b_pct], labels=[f"{ticker_a.upper()} {a_pct:.1f}%", f"{ticker_b.upper()} {b_pct:.1f}%"], autopct='%1.1f%%', startangle=140)
    axp.set_title("Robustness: share of random weights preferring each ticker")
    st.pyplot(figp)

# Price history plot
st.subheader("Price history")
fig, ax = plt.subplots(figsize=(10,4))
if not tA['history'].empty:
    ax.plot(tA['history'].index, tA['history']['Close'], label=ticker_a.upper())
if not tB['history'].empty:
    ax.plot(tB['history'].index, tB['history']['Close'], label=ticker_b.upper())
ax.set_title(f"{ticker_a.upper()} vs {ticker_b.upper()} Close Price")
ax.legend()
st.pyplot(fig)

# Heatmap: discount rate vs growth — show % upside for ticker A relative to market price
st.subheader("Heatmap: Discount rate vs Growth (DCF sensitivity)")
# Build grid
drates = np.linspace(max(0.01, discount_rate-0.05), discount_rate+0.05, 25)
grows = np.linspace(max(-0.02, default_fc_growth-0.05), default_fc_growth+0.10, 25)
heat = np.zeros((len(grows), len(drates)))
# compute intrinsic per share for grid for ticker A
for i,g in enumerate(grows):
    for j,d in enumerate(drates):
        dcf_tmp = dcf_5yr_simple(baseA, g, tgA, d, sharesA)
        if dcf_tmp and mA.get('regularMarketPrice'):
            heat[i,j] = (dcf_tmp['intrinsic_per_share'] - mA.get('regularMarketPrice')) / mA.get('regularMarketPrice') * 100.0
        else:
            heat[i,j] = np.nan

fighm, axhm = plt.subplots(figsize=(8,4))
im = axhm.imshow(heat, aspect='auto', origin='lower', extent=[drates.min(), drates.max(), grows.min(), grows.max()], cmap='RdYlGn')
axhm.set_xlabel("Discount rate")
axhm.set_ylabel("5yr growth rate")
axhm.set_title(f"{ticker_a.upper()} intrinsic vs market % (green positive upside)")
plt.colorbar(im, ax=axhm, label="% intrinsic - market")
st.pyplot(fighm)

# Allow downloads
st.subheader("Downloads: CSV / JSON / PDF / PPTX")
colD1, colD2, colD3, colD4 = st.columns([1,1,1,1])
combined_csv = combined_df.reset_index().to_csv(index=False).encode('utf-8')
score_csv = score_df.reset_index().to_csv(index=False).encode('utf-8')
dcf_csv = dcf_show.reset_index().to_csv(index=False).encode('utf-8')
with colD1:
    st.download_button("Download combined metrics (CSV)", data=combined_csv, file_name="combined_metrics.csv", mime="text/csv")
with colD2:
    st.download_button("Download scores (CSV)", data=score_csv, file_name="composite_scores.csv", mime="text/csv")
with colD3:
    st.download_button("Download dcf results (CSV)", data=dcf_csv, file_name="dcf_results.csv", mime="text/csv")
with colD4:
    # PDF
    if pdf_button:
        analysis_text = textwrap.dedent(f"""
            Comparative analysis generated: {datetime.datetime.utcnow().isoformat()} UTC
            Tickers: {ticker_a.upper()} vs {ticker_b.upper()}
            DCF assumptions: discount_rate={discount_rate}, default 5yr growth={default_fc_growth}
        """)
        pdfbuf = make_pdf_report(analysis_text, combined_df, {ticker_a.upper(): tA['history'], ticker_b.upper(): tB['history']}, sens_results)
        b64 = base64.b64encode(pdfbuf.read()).decode()
        href = f'<a href="data:application/pdf;base64,{b64}" download="coke_pepsi_report.pdf"> Download PDF report</a>'
        st.markdown(href, unsafe_allow_html=True)
    # PPTX
    if pptx_button:
        if not PPTX_AVAILABLE:
            st.error("python-pptx is not available. Install python-pptx to enable PPTX export.")
        else:
            title_text = f"Generated: {datetime.datetime.utcnow().isoformat()} — {ticker_a.upper()} vs {ticker_b.upper()}"
            dcf_results = {ticker_a.upper(): dcfA, ticker_b.upper(): dcfB}
            out_ppt = make_pptx_report(title_text, combined_df, dcf_results, {ticker_a.upper(): tA['history'], ticker_b.upper(): tB['history']}, out_path=f"{ticker_a}_{ticker_b}_summary.pptx")
            # serve as download
            with open(out_ppt, "rb") as f:
                data = f.read()
            b64 = base64.b64encode(data).decode()
            href = f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64}" download="{ticker_a}_{ticker_b}_summary.pptx"> Download PPTX Summary</a>'
            st.markdown(href, unsafe_allow_html=True)

# Show available snapshots for reproducibility
st.subheader("Stored snapshots (recent) — reproducibility")
sA_rows = get_snapshots_for_ticker(ticker_a, limit=5)
sB_rows = get_snapshots_for_ticker(ticker_b, limit=5)
def show_snapshots(rows, label):
    if not rows:
        st.write(f"No snapshots for {label}")
        return
    for r in rows:
        st.write(f"ID: {r[0]} | Ticker: {r[1]} | TS: {r[2]}")
with st.expander("Show recent snapshots for Ticker A"):
    show_snapshots(sA_rows, ticker_a)
with st.expander("Show recent snapshots for Ticker B"):
    show_snapshots(sB_rows, ticker_b)

# Footer & notes
st.markdown("---")
st.markdown("Notes: This tool is Designed by Randy Singh from KNet consulting & intended for real-world comparison and research. Discounted Cash Flow (DCF) here is a simplified template — real-world valuation requires full schedules (tax, capex, working capital, debt schedules). The app stores snapshots locally in `ticker_snapshots.db` for audit and reproducibility.")
